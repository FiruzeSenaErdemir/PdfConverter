import tabula
import pandas as pd
import fitz  # PyMuPDF kütüphanesi
from pptx import Presentation
from pptx.util import Inches
import os
from pdf2docx import parse

def extract_text_from_pdf(pdf_file_path):
    pdf_document = fitz.open(pdf_file_path)
    text = ''
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    pdf_document.close()
    return text

def pdf_to_excel(pdf_file_path, excel_file_path):
    # Pdfdeki tabloları okuma
    tables = tabula.read_pdf(pdf_file_path, pages='all')

    # Texti PDF dosyasından çekme
    text = extract_text_from_pdf(pdf_file_path)

    # Tabloları ayırma
    with pd.ExcelWriter(excel_file_path) as writer:
        for i, table in enumerate(tables):
            table.to_excel(writer, sheet_name=f'Table {i+1}')

        # Text adı altında excellde sheet oluşturup tablo dışında kalan datayı texte yazma
        pd.DataFrame({'Text': [text]}).to_excel(writer, sheet_name='Text', index=False)

def pdf_to_docx(file_path):

    pdf_file = file_path
    word_file = pdf_file + "-cevrilmis.docx"
    parse(pdf_file, word_file, start=0, end=None)

def pdf_to_pptx(file_path):
    # PDF dosyasının adı ve yolunu belirtin
    pdf_path = file_path

    # PowerPoint dosyasının adını belirtin
    pptx_path = pdf_path+'-cevrilmis.pptx'

    # Yeni bir PowerPoint sunusu oluşturma
    ppt = Presentation()

    # PDF dosyasını açma
    pdf_document = fitz.open(pdf_path)

    # Her bir PDF sayfasını PowerPoint slaydına ekleme
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_bytes = page.get_pixmap().tobytes()
        img_path = f"page_{page_num}.png"

        # Resmi diske kaydetme
        with open(img_path, "wb") as img_file:
            img_file.write(image_bytes)

        # Resmi PowerPoint slaydına ekleme
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # 5 numaralı layout genellikle boş bir sayfa sağlar
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

        # Diskteki resmi silme
        os.remove(img_path)

    # PowerPoint dosyasını kaydetme
    ppt.save(pptx_path)

    # PDF dosyasını kapatma
    pdf_document.close()

#kullanım
pdf_to_excel('örnek.pdf', 'örnek-excel-cevrilmis.xlsx')
pdf_to_pptx('örnek.pdf')
pdf_to_docx('örnek.pdf')
